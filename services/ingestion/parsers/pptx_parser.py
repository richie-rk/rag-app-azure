"""PPTX parser using python-pptx."""

from pptx import Presentation

from .base import BaseParser, ParsedPage


class PptxParser(BaseParser):
    def parse(self, file_path: str, file_name: str) -> list[ParsedPage]:
        """Parse PPTX into one ParsedPage per slide."""
        prs = Presentation(file_path)
        pages = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip(" |"):
                            texts.append(row_text)

            content = "\n".join(texts)
            if content:
                pages.append(
                    ParsedPage(
                        page_number=slide_num,
                        content=content,
                        metadata={"source": file_name, "page": slide_num},
                    )
                )

        return pages
